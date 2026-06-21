"""
ingest.py
Week 1: Load PDFs and PPTX files -> chunk text -> generate embeddings -> store in Qdrant (local, on-disk)

Run: python src/ingest.py
"""

import os
import uuid
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct

DATA_DIR = Path(os.environ.get("DATA_DIR", str(Path(__file__).parent.parent / "data")))
COLLECTION_NAME = "studymate_chunks"
CHUNK_SIZE = 500       # characters per chunk
CHUNK_OVERLAP = 100     # overlap between consecutive chunks
EMBEDDING_MODEL = "BAAI/bge-base-en-v1.5"  # good free local embedding model


def load_pdf_text(pdf_path: Path) -> str:
    """Extract raw text from a PDF file."""
    reader = PdfReader(str(pdf_path))
    text = ""
    for page in reader.pages:
        page_text = page.extract_text() or ""
        text += page_text + "\n"
    return text


def load_pptx_text(pptx_path: Path) -> str:
    """Extract raw text from a PPTX file (all text boxes, slide by slide)."""
    prs = Presentation(str(pptx_path))
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = f"\n[Slide {slide_num}]\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs)
                    if para_text.strip():
                        slide_text += para_text + "\n"
            # Also grab text from tables if present
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_text += cell.text + " "
                    slide_text += "\n"
        text += slide_text
    return text


def chunk_text(text: str, source: str, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """
    Split text into overlapping chunks.
    Each chunk keeps track of its source file for citation later.
    """
    chunks = []
    start = 0
    text = " ".join(text.split())  # normalize whitespace
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append({"text": chunk, "source": source})
        start += chunk_size - overlap
    return chunks


def main():
    print("Loading embedding model...")
    model = SentenceTransformer(EMBEDDING_MODEL)
    embedding_dim = model.get_sentence_embedding_dimension()

    print("Connecting to Qdrant (local, on-disk)...")
    client = QdrantClient(path="./qdrant_data")  # stores data in a local folder, no server needed

    # Recreate collection fresh each time you run ingestion (simple for now)
    client.recreate_collection(
        collection_name=COLLECTION_NAME,
        vectors_config=VectorParams(size=embedding_dim, distance=Distance.COSINE),
    )

    all_chunks = []
    pdf_files = list(DATA_DIR.glob("*.pdf"))
    pptx_files = list(DATA_DIR.glob("*.pptx"))

    if not pdf_files and not pptx_files:
        print(f"No PDF or PPTX files found in {DATA_DIR}. Add your subject files there and re-run.")
        print("Note: old .ppt files are NOT supported directly - convert to .pptx or .pdf first.")
        return

    for pdf_path in pdf_files:
        print(f"Processing {pdf_path.name}...")
        text = load_pdf_text(pdf_path)
        chunks = chunk_text(text, source=pdf_path.name)
        all_chunks.extend(chunks)
        print(f"  -> {len(chunks)} chunks")

    for pptx_path in pptx_files:
        print(f"Processing {pptx_path.name}...")
        text = load_pptx_text(pptx_path)
        chunks = chunk_text(text, source=pptx_path.name)
        all_chunks.extend(chunks)
        print(f"  -> {len(chunks)} chunks")

    print(f"\nTotal chunks across all files: {len(all_chunks)}")
    print("Generating embeddings...")

    texts = [c["text"] for c in all_chunks]
    embeddings = model.encode(texts, show_progress_bar=True, normalize_embeddings=True)

    print("Uploading to Qdrant...")
    points = [
        PointStruct(
            id=str(uuid.uuid4()),
            vector=embedding.tolist(),
            payload={"text": chunk["text"], "source": chunk["source"]},
        )
        for chunk, embedding in zip(all_chunks, embeddings)
    ]

    client.upsert(collection_name=COLLECTION_NAME, points=points)
    print(f"\nDone! {len(points)} chunks stored in collection '{COLLECTION_NAME}'.")


if __name__ == "__main__":
    main()